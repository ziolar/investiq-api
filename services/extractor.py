import io
from pathlib import Path


def extract_text(content: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(content)
    elif suffix in (".pptx", ".ppt"):
        return _extract_pptx(content)
    elif suffix in (".xlsx", ".xls"):
        return _extract_excel(content)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def _extract_pdf(content: bytes) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=content, filetype="pdf")
    parts = []
    for page in doc:
        parts.append(page.get_text())
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"[幻灯片 {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_excel(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line.replace("\t", ""):
                rows.append(line)
        if rows:
            parts.append(f"[工作表: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(parts)
